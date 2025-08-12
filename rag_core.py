import os, json, uuid
from typing import List, Tuple, Dict
import numpy as np
import faiss
from sentence_transformers import SentenceTransformer
from pypdf import PdfReader
from pptx import Presentation
import tiktoken

STORE_BASE = "store"
os.makedirs(STORE_BASE, exist_ok=True)

# Global objects (loaded once)
ENC = tiktoken.get_encoding("cl100k_base")
EMB = SentenceTransformer("sentence-transformers/all-MiniLM-L6-v2")


def tokenize_len(text: str) -> int:
    return len(ENC.encode(text or ""))


def chunk_text(text: str, max_tokens=350, overlap=60) -> List[str]:
    tokens = ENC.encode(text or "")
    chunks = []
    start = 0
    while start < len(tokens):
        end = min(start + max_tokens, len(tokens))
        chunk = ENC.decode(tokens[start:end])
        chunks.append(chunk)
        start = max(end - overlap, end)
    return chunks


def read_pdf(path: str) -> List[Tuple[str, str]]:
    reader = PdfReader(path)
    out = []
    for i, page in enumerate(reader.pages, start=1):
        txt = page.extract_text() or ""
        out.append((f"page_{i}", txt.strip()))
    return out


def read_pptx(path: str) -> List[Tuple[str, str]]:
    prs = Presentation(path)
    out = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
        out.append((f"slide_{i}", "\n".join(texts).strip()))
    return out


def embed_texts(texts: List[str]) -> np.ndarray:
    vecs = EMB.encode(texts, normalize_embeddings=True)
    return np.array(vecs, dtype="float32")


def build_index(chunks: List[Dict]):
    vecs = embed_texts([c["text"] for c in chunks])  # normalized
    dim = vecs.shape[1]
    index = faiss.IndexFlatIP(dim)
    index.add(vecs)

    # Topic centroid from longest chunks (proxy for doc theme)
    lengths = np.array([len(c["text"]) for c in chunks])
    top_idx = lengths.argsort()[::-1][: min(20, len(chunks))]
    centroid = vecs[top_idx].mean(axis=0, keepdims=True)
    centroid /= (np.linalg.norm(centroid) + 1e-12)
    return index, centroid


def save_store(base: str, doc_id: str, index, centroid: np.ndarray, chunks: List[Dict]):
    import json
    path = os.path.join(base, doc_id)
    os.makedirs(path, exist_ok=True)
    faiss.write_index(index, os.path.join(path, "index.faiss"))
    np.save(os.path.join(path, "centroid.npy"), centroid)
    with open(os.path.join(path, "chunks.jsonl"), "w", encoding="utf-8") as f:
        for c in chunks:
            f.write(json.dumps(c, ensure_ascii=False) + "\n")


def load_store(base: str, doc_id: str):
    path = os.path.join(base, doc_id)
    index = faiss.read_index(os.path.join(path, "index.faiss"))
    centroid = np.load(os.path.join(path, "centroid.npy"))
    chunks = []
    with open(os.path.join(path, "chunks.jsonl"), "r", encoding="utf-8") as f:
        for line in f:
            chunks.append(json.loads(line))
    return type("Doc", (), {"doc_id": doc_id, "index": index, "centroid": centroid, "chunks": chunks})


def ingest_file(path: str, store_base: str = STORE_BASE) -> str:
    name = path.lower()
    if name.endswith(".pdf"):
        units = read_pdf(path)
    elif name.endswith(".pptx") or name.endswith(".ppt"):
        units = read_pptx(path)
    else:
        raise ValueError("Unsupported file type (pdf, ppt, pptx)")

    chunks: List[Dict] = []
    for ref, text in units:
        if not text.strip():
            continue
        for i, ch in enumerate(chunk_text(text)):
            chunks.append({"id": f"{ref}_chunk_{i+1}", "ref": ref, "text": ch})

    index, centroid = build_index(chunks)
    doc_id = uuid.uuid4().hex
    save_store(store_base, doc_id, index, centroid, chunks)
    return doc_id


def is_on_topic(question: str, centroid: np.ndarray, on_topic_thresh=0.25):
    qv = embed_texts([question])
    sim = float(np.dot(qv[0], centroid[0]))
    return (sim >= on_topic_thresh, sim)


def retrieve(doc, question: str, k=5) -> List[Dict]:
    qv = embed_texts([question])
    D, I = doc.index.search(qv, k)
    out = []
    for score, idx in zip(D[0], I[0]):
        if idx == -1:
            continue
        out.append({**doc.chunks[idx], "score": float(score)})
    return out


def support_strength(hits: List[Dict]) -> float:
    if not hits:
        return 0.0
    top = sorted([h["score"] for h in hits], reverse=True)[:3]
    return float(sum(top) / len(top)) if top else 0.0